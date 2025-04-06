import os
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
import logging
from datetime import datetime
import zipfile
import json
import tempfile
import shutil
import warnings
from tqdm import tqdm
import cv2
import numpy as np
from PIL import Image
import pytesseract

# Системные константы
CURRENT_USER = "spy12297"
CURRENT_TIME = datetime.strptime("2025-04-06 17:33:49", "%Y-%m-%d %H:%M:%S")
CHUNK_SIZE = 1000
BATCH_SIZE = 50
UNRAR_PATH = r"C:\Program Files\UnRAR.exe"
TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Настройка Tesseract
pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

# Настройка игнорирования предупреждений
warnings.filterwarnings('ignore', category=UserWarning)

# Инициализация поддержки RAR
try:
    import rarfile
    rarfile.UNRAR_TOOL = UNRAR_PATH
    RARFILE_AVAILABLE = True
    print(f"RAR support enabled. Using UnRAR at: {UNRAR_PATH}")
except ImportError:
    RARFILE_AVAILABLE = False
    print("Warning: rarfile not installed. RAR archives will be skipped.")

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_search.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger('DocumentSearch')

class DocumentSearch:
    def __init__(self):
        self.base_path = Path("C:/Users/andre/ChromaDBGPTDeepSeek/Data")
        self.db_path = Path("C:/Users/andre/ChromaDBGPTDeepSeek/VectorDB")
        self.temp_path = Path("C:/Users/andre/ChromaDBGPTDeepSeek/Temp")
        self.current_user = CURRENT_USER
        self.last_index_time = CURRENT_TIME
        self.components = {}
        self.batch_size = BATCH_SIZE
        self.chunk_size = CHUNK_SIZE
        
        # Расширения файлов для OCR
        self.ocr_extensions = {'.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'}
        self.min_text_length = 50  # Минимальная длина текста для сохранения
        
        # Создаём необходимые директории
        for path in [self.base_path, self.db_path, self.temp_path]:
            path.mkdir(parents=True, exist_ok=True)
        
        # Инициализация компонентов
        self._initialize_components()

    def _initialize_components(self) -> None:
        """Инициализация всех необходимых компонентов"""
        self.components = {
            'chromadb': False,
            'docx': False,
            'xlsx': False,
            'pptx': False,
            'msg': False,
            'rar': RARFILE_AVAILABLE,
            'zip': True,
            'tesseract': self._check_tesseract()
        }
        
        # Инициализация ChromaDB
        try:
            import chromadb
            from chromadb.config import Settings
            
            settings = Settings(
                anonymized_telemetry=False,
                allow_reset=True,
                is_persistent=True
            )
            
            self.chroma_client = chromadb.PersistentClient(
                path=str(self.db_path),
                settings=settings
            )

            try:
                self.chroma_client.delete_collection("documents")
                logger.info("Deleted existing collection")
            except:
                logger.info("No existing collection to delete")
                
            self.collection = self.chroma_client.create_collection(
                name="documents",
                metadata={
                    "description": "Local document collection",
                    "created_by": self.current_user,
                    "created_at": self.last_index_time.isoformat()
                }
            )
            
            self.components['chromadb'] = True
            logger.info("ChromaDB initialized successfully")
        except Exception as e:
            logger.error(f"Error initializing ChromaDB: {e}")
            raise

        # Инициализация поддержки форматов документов
        self._initialize_document_support()

    def _check_tesseract(self) -> bool:
        """Проверка доступности Tesseract OCR"""
        try:
            version = pytesseract.get_tesseract_version()
            logger.info(f"Tesseract version: {version}")
            return True
        except Exception as e:
            logger.warning(f"Tesseract OCR not available: {e}")
            return False

    def _initialize_document_support(self) -> None:
        """Инициализация поддержки различных форматов документов"""
        # Word (.docx) support
        try:
            from docx import Document
            self.components['docx'] = True
            logger.info("Word (.docx) support enabled")
        except ImportError:
            logger.warning("Word (.docx) support disabled")

        # Excel (.xlsx) support
        try:
            import pandas as pd
            self.components['xlsx'] = True
            logger.info("Excel (.xlsx) support enabled")
        except ImportError:
            logger.warning("Excel (.xlsx) support disabled")

        # PowerPoint (.pptx) support
        try:
            from pptx import Presentation
            self.components['pptx'] = True
            logger.info("PowerPoint (.pptx) support enabled")
        except ImportError:
            logger.warning("PowerPoint (.pptx) support disabled")

        # Outlook (.msg) support
        try:
            import win32com.client
            import pythoncom
            self.components['msg'] = True
            logger.info("Outlook (.msg) support enabled")
        except ImportError:
            logger.warning("Outlook (.msg) support disabled")

    def process_with_ocr(self, file_path: Path) -> str:
        """Обработка файла с помощью OCR"""
        try:
            # Прямая обработка PDF и изображений через Tesseract
            text = pytesseract.image_to_string(
                str(file_path),
                lang='rus+eng',
                config='--psm 1 --oem 1'  # Использовать LSTM OCR Engine Mode
            )
            
            # Если текст слишком короткий или пустой, попробуем дополнительную обработку
            if len(text.strip()) < self.min_text_length:
                # Для изображений используем предварительную обработку
                if file_path.suffix.lower() in {'.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'}:
                    image = cv2.imread(str(file_path))
                    if image is not None:
                        # Конвертация в градации серого
                        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
                        
                        # Удаление шума
                        denoised = cv2.fastNlMeansDenoising(gray)
                        
                        # Адаптивная бинаризация
                        binary = cv2.adaptiveThreshold(
                            denoised, 255,
                            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                            cv2.THRESH_BINARY,
                            11, 2
                        )
                        
                        # Повторное распознавание с обработанным изображением
                        text = pytesseract.image_to_string(
                            binary,
                            lang='rus+eng',
                            config='--psm 1 --oem 1'
                        )
            
            # Очистка и нормализация текста
            text = ' '.join(text.split())
            return text if len(text) >= self.min_text_length else ""
            
        except Exception as e:
            logger.error(f"Error processing file with OCR {file_path}: {e}")
            return ""

    def extract_text(self, file_path: Path) -> str:
        """Извлечение текста из файла"""
        try:
            ext = file_path.suffix.lower()
            
            # OCR для поддерживаемых форматов
            if ext in self.ocr_extensions and self.components['tesseract']:
                return self.process_with_ocr(file_path)
            
            # Обработка остальных форматов
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
                    
            elif ext == '.docx' and self.components['docx']:
                from docx import Document
                doc = Document(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                
            elif ext == '.xlsx' and self.components['xlsx']:
                import pandas as pd
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore")
                    df = pd.read_excel(file_path)
                return df.to_string()
                
            elif ext == '.pptx' and self.components['pptx']:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return '\n'.join(text_runs)
                
            elif ext == '.msg' and self.components['msg']:
                import win32com.client
                import pythoncom
                pythoncom.CoInitialize()
                try:
                    outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
                    msg = outlook.OpenSharedItem(str(file_path))
                    text = f"Subject: {msg.Subject}\nFrom: {msg.SenderName}\nBody:\n{msg.Body}"
                    return text
                finally:
                    pythoncom.CoUninitialize()
            
            logger.warning(f"Unsupported or unhandled file type: {ext}")
            return ""
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""

    def get_supported_extensions(self) -> List[str]:
        """Получение списка поддерживаемых расширений"""
        extensions = ['.txt']
        if self.components['docx']: extensions.append('.docx')
        if self.components['xlsx']: extensions.append('.xlsx')
        if self.components['pptx']: extensions.append('.pptx')
        if self.components['msg']: extensions.append('.msg')
        if self.components['tesseract']: extensions.extend(list(self.ocr_extensions))
        return extensions

    def get_supported_archives(self) -> List[str]:
        """Получение списка поддерживаемых архивов"""
        archives = []
        if self.components['zip']: archives.append('.zip')
        if self.components['rar']: archives.append('.rar')
        return archives

    def extract_archive(self, archive_path: Path, temp_dir: Path) -> List[Path]:
        """Извлечение файлов из архива"""
        extracted_files: List[Path] = []
        try:
            if archive_path.suffix.lower() == '.zip':
                with zipfile.ZipFile(archive_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                    extracted_files = list(temp_dir.rglob('*'))
            elif archive_path.suffix.lower() == '.rar' and self.components['rar']:
                with rarfile.RarFile(archive_path, 'r') as rar_ref:
                    rar_ref.extractall(temp_dir)
                    extracted_files = list(temp_dir.rglob('*'))
            
            return [f for f in extracted_files if f.is_file()]
        except Exception as e:
            logger.error(f"Error extracting archive {archive_path}: {e}")
            return []

    def _batch_add_to_collection(self, documents: List[str], ids: List[str], metadatas: List[Dict[str, Any]]) -> None:
        """Добавление документов в коллекцию батчами"""
        if not documents:
            return
            
        try:
            self.collection.add(
                documents=documents,
                ids=ids,
                metadatas=metadatas
            )
        except Exception as e:
            logger.error(f"Error adding batch to collection: {e}")
            # Пробуем добавить по одному документу
            for i in range(len(documents)):
                try:
                    self.collection.add(
                        documents=[documents[i]],
                        ids=[ids[i]],
                        metadatas=[metadatas[i]]
                    )
                except Exception as e:
                    logger.error(f"Error adding document {ids[i]}: {e}")

    def process_archive(self, archive_path: Path, stats: Dict[str, Any]) -> None:
        """Обработка архива и его содержимого"""
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            extracted_files = self.extract_archive(archive_path, temp_path)
            
            for file_path in extracted_files:
                if file_path.suffix.lower() in self.get_supported_extensions():
                    try:
                        text = self.extract_text(file_path)
                        if text:
                            # Создаем специальный путь для файла из архива
                            archive_file_path = f"{archive_path}::{file_path.name}"
                            
                            # Разбиваем текст на части
                            chunks = [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]
                            
                            for i, chunk in enumerate(chunks):
                                self.current_batch['documents'].append(chunk)
                                self.current_batch['ids'].append(f"{archive_file_path}_{i}")
                                self.current_batch['metadatas'].append({
                                    "filename": file_path.name,
                                    "archive_path": str(archive_path),
                                    "path": archive_file_path,
                                    "extension": file_path.suffix,
                                    "chunk": i,
                                    "total_chunks": len(chunks),
                                    "is_archived": True,
                                    "archive_name": archive_path.name,
                                    "indexed_at": self.last_index_time.isoformat(),
                                    "indexed_by": self.current_user
                                })
                            
                            stats['indexed_files'] += 1
                            
                            # Проверяем размер батча
                            if len(self.current_batch['documents']) >= self.batch_size:
                                self._batch_add_to_collection(
                                    self.current_batch['documents'],
                                    self.current_batch['ids'],
                                    self.current_batch['metadatas']
                                )
                                self.current_batch = {'documents': [], 'ids': [], 'metadatas': []}
                        
                    except Exception as e:
                        error_msg = f"Error processing file {file_path} from archive {archive_path}: {str(e)}"
                        logger.error(error_msg)
                        stats['errors'].append(error_msg)
                        stats['skipped_files'] += 1

    def index_documents(self) -> Dict[str, Any]:
        """Индексация документов"""
        self.last_index_time = datetime.utcnow()
        self.current_batch = {'documents': [], 'ids': [], 'metadatas': []}
        
        stats = {
            "processed_files": 0,
            "indexed_files": 0,
            "skipped_files": 0,
            "processed_archives": 0,
            "errors": [],
            "start_time": self.last_index_time.isoformat(),
            "supported_extensions": self.get_supported_extensions(),
            "supported_archives": self.get_supported_archives(),
            "user": self.current_user
        }
        
        try:
            files = list(self.base_path.rglob('*'))
            total_files = len([f for f in files if f.is_file()])
            
            logger.info(f"Found {total_files} total files to process")
            
            with tqdm(total=total_files, desc="Indexing files") as pbar:
                for file_path in files:
                    if not file_path.is_file():
                        continue
                    
                    stats['processed_files'] += 1
                    pbar.update(1)
                    
                    try:
                        # Проверяем, является ли файл архивом
                        if file_path.suffix.lower() in self.get_supported_archives():
                            stats['processed_archives'] += 1
                            self.process_archive(file_path, stats)
                            continue
                        
                        # Обработка обычных файлов
                        if file_path.suffix.lower() in self.get_supported_extensions():
                            text = self.extract_text(file_path)
                            if text:
                                chunks = [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]
                                
                                for i, chunk in enumerate(chunks):
                                    self.current_batch['documents'].append(chunk)
                                    self.current_batch['ids'].append(f"{file_path}_{i}")
                                    self.current_batch['metadatas'].append({
                                        "filename": file_path.name,
                                        "path": str(file_path),
                                        "extension": file_path.suffix,
                                        "chunk": i,
                                        "total_chunks": len(chunks),
                                        "is_archived": False,
                                        "indexed_at": self.last_index_time.isoformat(),
                                        "indexed_by": self.current_user
                                    })
                                
                                stats['indexed_files'] += 1
                                
                                if len(self.current_batch['documents']) >= self.batch_size:
                                    self._batch_add_to_collection(
                                        self.current_batch['documents'],
                                        self.current_batch['ids'],
                                        self.current_batch['metadatas']
                                    )
                                    self.current_batch = {'documents': [], 'ids': [], 'metadatas': []}
                            else:
                                stats['skipped_files'] += 1
                        
                    except Exception as e:
                        error_msg = f"Error processing {file_path}: {str(e)}"
                        logger.error(error_msg)
                        stats['errors'].append(error_msg)
                        stats['skipped_files'] += 1
            
            # Сохраняем оставшиеся документы
            if self.current_batch['documents']:
                self._batch_add_to_collection(
                    self.current_batch['documents'],
                    self.current_batch['ids'],
                    self.current_batch['metadatas']
                )
            
            stats['end_time'] = datetime.utcnow().isoformat()
            stats['duration'] = (datetime.fromisoformat(stats['end_time']) - 
                               datetime.fromisoformat(stats['start_time'])).total_seconds()
            
            logger.info(f"Indexing completed. Stats: {json.dumps(stats, indent=2)}")
            return stats
            
        except Exception as e:
            error_msg = f"Fatal error during indexing: {str(e)}"
            logger.error(error_msg)
            stats['errors'].append(error_msg)
            stats['end_time'] = datetime.utcnow().isoformat()
            return stats

    def search(self, query: str, n_results: int = 5) -> List[Dict[str, Any]]:
        """Поиск документов"""
        try:
            results = self.collection.query(
                query_texts=[query],
                n_results=n_results * 3,
                include=["documents", "metadatas", "distances"]
            )
            
            # Группируем результаты по файлам
            grouped_results = {}
            for i in range(len(results['ids'][0])):
                file_path = results['metadatas'][0][i]['path']
                current_score = 1 - results['distances'][0][i]
                
                if file_path not in grouped_results or current_score > grouped_results[file_path]['relevance_score']:
                    grouped_results[file_path] = {
                        'file_path': file_path,
                        'file_name': results['metadatas'][0][i]['filename'],
                        'relevance_score': current_score,
                        'snippet': results['documents'][0][i],
                        'metadata': results['metadatas'][0][i]
                    }
            
            # Сортируем и возвращаем топ n_results
            return sorted(
                grouped_results.values(),
                key=lambda x: x['relevance_score'],
                reverse=True
            )[:n_results]
            
        except Exception as e:
            logger.error(f"Error during search: {e}")
            return []

    def create_results_archive(self, results: List[Dict[str, Any]]) -> str:
        """Создание архива с найденными документами"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            archive_path = self.temp_path / f"search_results_{timestamp}.zip"
            
            with zipfile.ZipFile(archive_path, 'w') as zf:
                for result in results:
                    metadata = result.get('metadata', {})
                    is_archived = metadata.get('is_archived', False)
                    
                    if is_archived:
                        # Извлекаем файл из архива во временную директорию
                        archive_path = Path(metadata['archive_path'])
                        if archive_path.exists():
                            with tempfile.TemporaryDirectory() as temp_dir:
                                temp_path = Path(temp_dir)
                                if archive_path.suffix.lower() == '.zip':
                                    with zipfile.ZipFile(archive_path, 'r') as archive:
                                        archive.extract(metadata['filename'], temp_path)
                                elif archive_path.suffix.lower() == '.rar' and self.components['rar']:
                                    with rarfile.RarFile(archive_path, 'r') as archive:
                                        archive.extract(metadata['filename'], temp_path)
                                
                                extracted_file = temp_path / metadata['filename']
                                if extracted_file.exists():
                                    zf.write(extracted_file, arcname=f"{archive_path.stem}/{metadata['filename']}")
                    else:
                        # Обычный файл
                        source_path = Path(result['file_path'])
                        if source_path.exists():
                            zf.write(source_path, arcname=source_path.name)
            
            return str(archive_path)
            
        except Exception as e:
            logger.error(f"Error creating archive: {e}")
            return ""